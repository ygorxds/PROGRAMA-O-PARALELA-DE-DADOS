import pandas as pd
from concurrent.futures import ThreadPoolExecutor
import numpy as np
import matplotlib.pyplot as plt
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches
import time

# Leitura do arquivo Excel
def read_excel(file_path):
    df = pd.read_excel(file_path)
    return df

# Contagem de ocorrências de eventos específicos
def count_occurrences(df, column):
    return df[column].value_counts().to_dict()

# Cálculo de estatísticas
def calculate_statistics(df, column):
    mean = df[column].mean()
    median = df[column].median()
    stddev = df[column].std()
    return {'mean': mean, 'median': median, 'stddev': stddev}

# Transformações de dados (filtragem)
def transform_data(df, column, filter_value):
    if np.issubdtype(df[column].dtype, np.datetime64):
        return df[df[column] > filter_value]
    else:
        return df[df[column] > filter_value]

# Função para contagem de ocorrências em paralelo
def parallel_count_occurrences(df, column):
    with ThreadPoolExecutor() as executor:
        chunks = np.array_split(df, 4)
        results = executor.map(lambda chunk: count_occurrences(chunk, column), chunks)
    final_result = {}
    for result in results:
        for key, value in result.items():
            if key in final_result:
                final_result[key] += value
            else:
                final_result[key] = value
    return final_result

# Geração do relatório em PDF
def generate_report(statistics, output_path):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    pdf.cell(200, 10, txt="Relatório de Análise de Dados", ln=True, align='C')
    pdf.cell(200, 10, txt=f"Média: {statistics['mean']}", ln=True)
    pdf.cell(200, 10, txt=f"Mediana: {statistics['median']}", ln=True)
    pdf.cell(200, 10, txt=f"Desvio-Padrão: {statistics['stddev']}", ln=True)
    pdf.output(output_path)

# Geração da apresentação de slides
def generate_slides(statistics, output_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Relatório de Análise de Dados"
    content = slide.placeholders[1]
    content.text = f"Média: {statistics['mean']}\nMediana: {statistics['median']}\nDesvio-Padrão: {statistics['stddev']}"
    prs.save(output_path)

# Função principal
def main():
    start_time = time.time()  # Início da medição do tempo
    file_path = r'C:\Users\Ygor\Desktop\Shopping Barra.xlsx'
    try:
        df = read_excel(file_path)

        # Exemplo de contagem de ocorrências na coluna 1
        occurrences = parallel_count_occurrences(df, df.columns[0])
        print('Occurrences:', occurrences)

        # Exemplo de cálculo de estatísticas na coluna 2
        statistics = calculate_statistics(df, df.columns[1])
        print('Statistics:', statistics)

        # Exemplo de transformação de dados
        if np.issubdtype(df[df.columns[1]].dtype, np.datetime64):
            filter_value = pd.Timestamp('2019-12-11')
        else:
            filter_value = 10
        
        filtered_data = transform_data(df, df.columns[1], filter_value)
        print('Filtered Data:', filtered_data)

        # Gerar relatório e slides
        generate_report(statistics, 'report.pdf')
        generate_slides(statistics, 'presentation.pptx')
    except FileNotFoundError:
        print(f"Erro: O arquivo {file_path} não foi encontrado.")
    except Exception as e:
        print(f"Erro ao processar o arquivo: {e}")
    finally:
        end_time = time.time()  # Fim da medição do tempo
        execution_time = (end_time - start_time) * 1000  # Tempo de execução em milissegundos
        print(f"Tempo de execução: {execution_time:.2f} ms")

if __name__ == "__main__":
    main()
